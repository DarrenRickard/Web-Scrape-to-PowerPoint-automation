from pptx import Presentation
from pptx.util import Inches
from app import scrape_data, get_slide_count
import requests
from io import BytesIO

PRS_NAME = 'Presentation.pptx'
SLD_LAYOUT_TITLE_AND_CONTENT = 1
SLD_LAYOUT_TITLE_PIC_TEXT = 8
LOGO_URL = 'https://upload.wikimedia.org/wikipedia/commons/thumb/a/a2/AAFES_Redesigned_Logo_2011-vector.svg/2560px-AAFES_Redesigned_Logo_2011-vector.svg.png'

slide_count = get_slide_count()
for i in range(slide_count):
    try:
        itemTitle, imageURL, itemDescList, dimList, CRC = scrape_data()

        def download_image(url):
            response = requests.get(url)
            image_data = BytesIO(response.content)
            return image_data

        try:
            prs = Presentation(PRS_NAME)
        except Exception:
            prs = Presentation()
            
        # Change to standard vertical slide
        prs.slide_width = Inches(7.5)
        prs.slide_height = Inches(10)

        # Add slide
        slide_layout = prs.slide_layouts[SLD_LAYOUT_TITLE_PIC_TEXT]
        slide = prs.slides.add_slide(slide_layout)

        # Add title text
        title_placeholder = slide.shapes.title
        title_placeholder.text = f'{itemTitle}\nCRC # {CRC}\n$'
        title_placeholder.left = Inches(1.25)
        title_placeholder.top = Inches(0.5)
        title_placeholder.width = Inches(6)
        title_placeholder.height = Inches(2)

        # Add picture
        picture_placeholder = slide.placeholders[1]
        image = download_image(imageURL)
        pic = picture_placeholder.insert_picture(image)
        pic.width = Inches(5)
        pic.height = Inches(5)
        pic.left = Inches(1.25)
        pic.top = Inches(2.5)

        # Add body text
        text_placeholder = slide.placeholders[2]
        text_placeholder.left = Inches(1.25)
        text_placeholder.top = Inches(7.5)
        text_placeholder.width = Inches(6)
        text_placeholder.height = Inches(2)
        for item in itemDescList + dimList:
            current_text = text_placeholder.text
            if current_text: # if has a value already
                final_text = f'{current_text}\n{item}'
            else:
                final_text = item
            text_placeholder.text = final_text

        # for shape in slide.placeholders:
        #     print('%d %s' % (shape.placeholder_format.idx, shape.name))
        #     print(shape)

        prs.save(PRS_NAME)
    except Exception:
        print("Item not found. Skipping...")
        continue
    
print(f"{PRS_NAME} created!")